import os, pickle, io, base64, logging
import pandas as pd
from pptx import Presentation
from PyPDF2 import PdfReader
from docx import Document
from googleapiclient.discovery import build
from google.auth.transport.requests import Request
from googleapiclient.http import MediaIoBaseUpload, MediaIoBaseDownload
from reportlab.pdfgen import canvas
from typing import Dict, Any, Optional

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

EXTENSION_TO_MIME = {
    ".txt": "text/plain",
    ".csv": "text/csv",
    ".json": "application/json",
    ".pdf": "application/pdf",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".mp4": "video/mp4",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".py": "text/plain",
    ".ipynb": "application/json",
}

class GoogleDriveAPIClient:
    def __init__(self, user_id: str = "default", token_dir: str = "tokens"):
        self.user_id = user_id
        self.token_dir = token_dir
        self.token_file = os.path.join(token_dir, f"{user_id}_token.pickle")
        self.service = None
        self.scopes = [
            'https://www.googleapis.com/auth/drive',
            'https://www.googleapis.com/auth/drive.file'
        ]

    def authenticate(self) -> bool:
        """Authenticate Google Drive for given user_id"""
        try:
            if not os.path.exists(self.token_file):
                logger.error(f"No token found for {self.user_id}. Run auth_setup.py first.")
                return False

            with open(self.token_file, "rb") as token:
                creds = pickle.load(token)

            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
                with open(self.token_file, "wb") as token_out:
                    pickle.dump(creds, token_out)

            if not creds or not creds.valid:
                logger.error("Invalid credentials. Re-run OAuth for user.")
                return False

            self.service = build("drive", "v3", credentials=creds)
            logger.info(f"✅ Authenticated Google Drive for {self.user_id}")
            return True
        except Exception as e:
            logger.error(f"Auth failed for {self.user_id}: {e}")
            return False

    # 🔽 (keep your create_folder, list_directory, read_file, write_file, etc. methods as before)
    # Just ensure all methods use `self.service`


    def _get_mime_type(self, filename: str) -> str:
        """Infer MIME type from extension"""
        ext = os.path.splitext(filename)[1].lower()
        return EXTENSION_TO_MIME.get(ext, "application/octet-stream")

    def _make_pdf_bytes(self, text: str) -> bytes:
        """Generate a valid PDF file with given text"""
        buffer = io.BytesIO()
        c = canvas.Canvas(buffer)
        c.setFont("Helvetica", 12)
        c.drawString(100, 750, text)
        c.save()
        buffer.seek(0)
        return buffer.read()

    def create_folder(self, name: str, parent_id: str = None) -> Dict[str, Any]:
        try:
            if not self.service:
                return {"status": "error", "message": "Not authenticated"}

            folder_metadata = {
                'name': name,
                'mimeType': 'application/vnd.google-apps.folder'
            }
            if parent_id and parent_id != "root":
                folder_metadata['parents'] = [parent_id]

            folder = self.service.files().create(
                body=folder_metadata,
                fields='id,name,parents,createdTime,webViewLink'
            ).execute()

            return {"status": "success", "data": folder}
        except Exception as e:
            return {"status": "error", "message": f"Failed to create folder: {e}"}

    def list_directory(self, folder_id: str = None, max_results: int = 100) -> Dict[str, Any]:
        try:
            if not self.service:
                return {"status": "error", "message": "Not authenticated"}

            query = f"'{folder_id or 'root'}' in parents and trashed=false"
            results = self.service.files().list(
                q=query,
                pageSize=max_results,
                fields="files(id,name,mimeType,size,modifiedTime,webViewLink)"
            ).execute()

            return {"status": "success", "data": results.get("files", [])}
        except Exception as e:
            return {"status": "error", "message": f"Failed to list directory: {e}"}

    def navigate_path(self, path: str) -> Dict[str, Any]:
        try:
            if not self.service:
                return {"status": "error", "message": "Not authenticated"}

            if path in ("/", ""):
                return self.list_directory("root")

            parts = [p for p in path.split("/") if p]
            current_folder_id = "root"
            for part in parts:
                query = f"'{current_folder_id}' in parents and name='{part}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
                results = self.service.files().list(q=query, fields="files(id,name)").execute()
                items = results.get("files", [])
                if not items:
                    return {"status": "error", "message": f"Path '{path}' not found"}
                current_folder_id = items[0]['id']
            return self.list_directory(current_folder_id)
        except Exception as e:
            return {"status": "error", "message": f"Failed to navigate path: {e}"}

    def _get_file_id_by_name(self, name: str, parent_id: str = "root") -> Optional[str]:
        """Helper: find file ID by name in a given folder"""
        try:
            query = f"'{parent_id}' in parents and name='{name}' and trashed=false"
            results = self.service.files().list(
                q=query,
                fields="files(id,name)"
            ).execute()
            files = results.get("files", [])
            if not files:
                return None
            return files[0]["id"]
        except Exception as e:
            logger.error(f"Error finding file by name: {e}")
            return None

    def read_file(self, file_name_or_id: str, encoding: str = "utf-8", parent_id: str = "root") -> Dict[str, Any]:
        """
        Read file by name or ID
        Returns readable text if possible; otherwise Base64
        """
        try:
            if not self.service:
                return {"status": "error", "message": "Not authenticated"}

            # Resolve ID
            if len(file_name_or_id) > 20 and not file_name_or_id.endswith(tuple(EXTENSION_TO_MIME.keys())):
                file_id = file_name_or_id
            else:
                file_id = self._get_file_id_by_name(file_name_or_id, parent_id)
                if not file_id:
                    return {"status": "error", "message": f"File '{file_name_or_id}' not found"}

            metadata = self.service.files().get(
                fileId=file_id,
                fields='id,name,mimeType,size,modifiedTime'
            ).execute()

            request = self.service.files().get_media(fileId=file_id)
            file_io = io.BytesIO()
            downloader = MediaIoBaseDownload(file_io, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()

            file_io.seek(0)
            filename = metadata.get("name")
            ext = os.path.splitext(filename)[1].lower()

            # Handle different formats
            if ext in [".txt", ".py", ".csv", ".json", ".ipynb"]:
                content = file_io.getvalue().decode(encoding)
            elif ext == ".pdf":
                reader = PdfReader(file_io)
                content = "\n".join([page.extract_text() or "" for page in reader.pages])
            elif ext == ".xlsx":
                df = pd.read_excel(file_io)
                content = df.to_csv(index=False)
            elif ext == ".pptx":
                prs = Presentation(file_io)
                content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            else:
                # Binary → return Base64
                content = base64.b64encode(file_io.getvalue()).decode("utf-8")

            return {"status": "success", "data": {"metadata": metadata, "content": content}}

        except Exception as e:
            return {"status": "error", "message": f"Failed to read file: {e}"}

    def write_file(self, name: str, content: str, file_id: Optional[str] = None, parent_id: str = None) -> Dict[str, Any]:
        """Write file to Drive (supports text, PDF, Excel, PPTX, DOCX, binary)"""
        try:
            if not self.service:
                return {"status": "error", "message": "Not authenticated"}

            mime_type = self._get_mime_type(name)

            # PDF
            if mime_type == "application/pdf":
                file_bytes = self._make_pdf_bytes(content)
                file_stream = io.BytesIO(file_bytes)

            # Excel (.xlsx)
            elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                try:
                    file_bytes = base64.b64decode(content)
                    file_stream = io.BytesIO(file_bytes)
                except Exception:
                    df = pd.DataFrame([row.split(",") for row in content.splitlines()])
                    file_stream = io.BytesIO()
                    with pd.ExcelWriter(file_stream, engine="xlsxwriter") as writer:
                        df.to_excel(writer, index=False, header=False)
                    file_stream.seek(0)

            # PowerPoint (.pptx)
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                try:
                    file_bytes = base64.b64decode(content)
                    file_stream = io.BytesIO(file_bytes)
                except Exception:
                    prs = Presentation()
                    for line in content.splitlines():
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        slide.shapes.title.text = "Slide"
                        slide.placeholders[1].text = line
                    file_stream = io.BytesIO()
                    prs.save(file_stream)
                    file_stream.seek(0)

            # Word (.docx)
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                try:
                    file_bytes = base64.b64decode(content)
                    file_stream = io.BytesIO(file_bytes)
                except Exception:
                    doc = Document()
                    for line in content.splitlines():
                        doc.add_paragraph(line)
                    file_stream = io.BytesIO()
                    doc.save(file_stream)
                    file_stream.seek(0)

            # Plain text / JSON
            elif mime_type.startswith("text/") or mime_type == "application/json":
                file_stream = io.BytesIO(content.encode("utf-8"))

            # Fallback → binary/base64
            else:
                try:
                    file_bytes = base64.b64decode(content)
                except Exception:
                    file_bytes = content.encode("utf-8")
                file_stream = io.BytesIO(file_bytes)

            # Upload to Drive
            file_stream.seek(0)
            media = MediaIoBaseUpload(file_stream, mimetype=mime_type, resumable=True)

            if file_id:
                updated = self.service.files().update(
                    fileId=file_id,
                    body={'name': name},
                    media_body=media,
                    fields='id,name,size,modifiedTime,webViewLink'
                ).execute()
                return {"status": "success", "data": updated}
            else:
                metadata = {'name': name}
                if parent_id and parent_id != "root":
                    metadata['parents'] = [parent_id]
                created = self.service.files().create(
                    body=metadata,
                    media_body=media,
                    fields='id,name,size,createdTime,webViewLink,parents'
                ).execute()
                return {"status": "success", "data": created}

        except Exception as e:
            return {"status": "error", "message": f"Failed to write file: {e}"}

